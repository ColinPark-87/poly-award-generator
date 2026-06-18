# -*- coding: utf-8 -*-
"""
분당엠폴리 'Best Essay' (Certificate of Excellence) 전용 상장 템플릿 베이킹.

원본 PPT(Raw/0618/26_Best essay 1.pptx, 35슬라이드=학생35명)에서 슬라이드 1장만 떼어
변수(반코드·이름)를 비운 깨끗한 1쪽 PDF 템플릿을 만든다. 정적 텍스트(제목 'Best Essay',
'Certificate of Excellence', 서명 'Bobin Park' 등)는 PPT에 임베드된 폰트(Zen Old Mincho 등)
그대로 PowerPoint COM PDF 내보내기로 베이킹된다. 런타임(generator.py)은 비워둔 자리에
학생별 반코드·이름만 올린다.

추가 편집(사용자 요청): 'Bundang MPOLY' 의 'M'(MPOLY 첫 글자)을 Poly 로열블루(#004EA2)로.
  → 본문 TextBox13("...at Bundang MPOLY")·서명 TextBox14("Bundang MPOLY Director")에서
    'MPOLY' 의 'M' 한 글자만 run 분리하여 파란색 지정.

산출물:
  templates/분당엠폴리/best_essay.pdf            (베이킹된 빈 템플릿)
  test_output/bestessay_ref.pdf                  (원본 슬라이드1, 좌표/육안 참조용)
  test_output/bestessay_ref_coords.json          (반코드·이름 placeholder 좌표)
  test_output/bestessay_template_preview.png

재생성: python build_bestessay_template.py   (Windows + PowerPoint COM 필요)
"""
import os
import json
import shutil
import tempfile

import config

SRC_PPTX = r"C:/Users/user/Desktop/Colin 작업폴더/Raw/0618/26_Best essay 1.pptx"
DST_PDF  = os.path.join(config.TEMPLATE_DIR, "분당엠폴리", "best_essay.pdf")
PREVIEW_DIR = os.path.join(os.path.dirname(__file__), "test_output")

POLY_BLUE = (0x00, 0x4E, 0xA2)   # Royal Blue (Poly 브랜드)

# PPT 텍스트박스 이름 (검증된 슬라이드 구조)
TB_CLASS = "TextBox 11"   # 반코드 (예: 5HO1_1)
TB_NAME  = "TextBox 12"   # 이름   (예: 정현우 (Daniel Jung))
TB_BODY  = "TextBox 13"   # "...at Bundang MPOLY"
TB_SIG   = "TextBox 14"   # "Bundang MPOLY Director"


def _rgb_int(r, g, b):
    """python-pptx RGBColor 용 0xRRGGBB int."""
    return (r << 16) | (g << 8) | b


def _blue_the_M(text_frame):
    """text_frame 안 'MPOLY' 의 첫 글자 'M' run 을 분리해 파란색 지정."""
    from pptx.dml.color import RGBColor
    for para in text_frame.paragraphs:
        full = "".join(r.text for r in para.runs)
        idx = full.find("MPOLY")
        if idx < 0:
            continue
        # 단순화: para 를 'M' 앞 / 'M' / 나머지 3 구간으로 재구성하되 첫 run 서식을 복제
        if not para.runs:
            continue
        base = para.runs[0]
        font_name = base.font.name
        font_size = base.font.size
        bold = base.font.bold
        try:
            base_color = base.font.color.rgb
        except Exception:
            base_color = None

        before = full[:idx]
        rest   = full[idx + 1:]   # 'POLY...' 이후 (M 다음 전체)

        # 기존 run 모두 제거
        for r in list(para.runs):
            r._r.getparent().remove(r._r)

        def _add(txt, blue=False):
            run = para.add_run()
            run.text = txt
            if font_name:
                run.font.name = font_name
            if font_size:
                run.font.size = font_size
            if bold is not None:
                run.font.bold = bold
            if blue:
                run.font.color.rgb = RGBColor(*POLY_BLUE)
            elif base_color is not None:
                run.font.color.rgb = base_color

        if before:
            _add(before)
        _add("M", blue=True)
        if rest:
            _add(rest)
        return True
    return False


def _export_pdf_via_com(pptx_path, pdf_path):
    """PowerPoint COM 으로 .pptx → .pdf (임베드 폰트 사용)."""
    import win32com.client
    pp = win32com.client.Dispatch("PowerPoint.Application")
    try:
        deck = pp.Presentations.Open(os.path.abspath(pptx_path),
                                     WithWindow=False, ReadOnly=True)
        deck.SaveAs(os.path.abspath(pdf_path), 32)   # 32 = ppSaveAsPDF
        deck.Close()
    finally:
        pp.Quit()


def _read_placeholder_coords(ref_pdf):
    """원본 슬라이드1 PDF 에서 반코드/이름 placeholder span 좌표를 읽는다."""
    import fitz
    doc = fitz.open(ref_pdf)
    page = doc[0]
    spans = [s for b in page.get_text("dict")["blocks"] if b["type"] == 0
             for l in b["lines"] for s in l["spans"]]
    out = {"page_w": page.rect.width, "page_h": page.rect.height}

    def _pack(s):
        c = s["color"]
        return {"bbox": list(s["bbox"]), "origin": list(s["origin"]),
                "size": s["size"], "color": [(c >> 16) & 255, (c >> 8) & 255, c & 255]}

    # 반코드: '5HO1_1' (또는 _ 포함 영숫자 코드, y 0.45~0.6 영역)
    ph = config.page_h = page.rect.height
    cls_s = next((s for s in spans if s["text"].strip().rstrip("_").replace("_", "").isalnum()
                  and "_" in s["text"] and ph * 0.46 < s["origin"][1] < ph * 0.60), None)
    name_s = next((s for s in spans if "(" in s["text"] and ")" in s["text"]
                   and ph * 0.55 < s["origin"][1] < ph * 0.72), None)
    if cls_s is None or name_s is None:
        # 폴백: 가장 가까운 텍스트
        raise RuntimeError(f"placeholder 미검출 cls={cls_s is not None} name={name_s is not None}")
    out["class"] = _pack(cls_s)
    out["name"]  = _pack(name_s)
    doc.close()
    return out


def _blank_and_blue(pptx_path, out_pptx):
    """슬라이드1만 남기고 반코드·이름 비우기 + MPOLY 의 M 파란색."""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    # 슬라이드 1장만 남기기: 뒤에서부터 제거
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for sld in slides[1:]:
        xml_slides.remove(sld)
    slide = prs.slides[0]
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if sh.name in (TB_CLASS, TB_NAME):
            # 텍스트만 비우기 (서식·박스 위치 유지 위해 첫 run 텍스트 ''로)
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    r.text = ""
        elif sh.name in (TB_BODY, TB_SIG):
            _blue_the_M(sh.text_frame)
    prs.save(out_pptx)


def main():
    if not os.path.exists(SRC_PPTX):
        raise FileNotFoundError(SRC_PPTX)
    os.makedirs(PREVIEW_DIR, exist_ok=True)
    os.makedirs(os.path.dirname(DST_PDF), exist_ok=True)

    tmpd = tempfile.mkdtemp()
    ref_pptx   = os.path.join(tmpd, "ref_slide1.pptx")
    blank_pptx = os.path.join(tmpd, "blank_slide1.pptx")
    ref_pdf    = os.path.join(PREVIEW_DIR, "bestessay_ref.pdf")

    # 1) 슬라이드1만 남긴 원본(좌표 참조용)
    from pptx import Presentation
    prs = Presentation(SRC_PPTX)
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides)[1:]:
        xml_slides.remove(sld)
    prs.save(ref_pptx)
    _export_pdf_via_com(ref_pptx, ref_pdf)
    coords = _read_placeholder_coords(ref_pdf)
    with open(os.path.join(PREVIEW_DIR, "bestessay_ref_coords.json"), "w", encoding="utf-8") as f:
        json.dump(coords, f, ensure_ascii=False, indent=2)
    print("[coords] class:", coords["class"])
    print("[coords] name :", coords["name"])

    # 2) 변수 비우기 + MPOLY M 파란색 → 템플릿 PDF
    _blank_and_blue(SRC_PPTX, blank_pptx)
    _export_pdf_via_com(blank_pptx, DST_PDF)
    print("[OK] template ->", DST_PDF)

    # 3) 프리뷰
    import fitz
    from PIL import Image
    doc = fitz.open(DST_PDF)
    pix = doc[0].get_pixmap(matrix=fitz.Matrix(150 / 72, 150 / 72))
    Image.frombytes("RGB", [pix.width, pix.height], pix.samples).save(
        os.path.join(PREVIEW_DIR, "bestessay_template_preview.png"))
    doc.close()
    shutil.rmtree(tmpd, ignore_errors=True)
    print("완료. 프리뷰:", os.path.join(PREVIEW_DIR, "bestessay_template_preview.png"))


if __name__ == "__main__":
    main()
